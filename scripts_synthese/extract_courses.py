#!/usr/bin/env python3
"""
extract_courses.py — Extraction et nettoyage de tous les cours
Utilisé par les agents Claude Code avant la phase de synthèse.
"""

import pdfplumber
import subprocess
import openpyxl
import os
import re
import json
from pathlib import Path

def extract_and_clean_pdf(filepath):
    """Extraire et nettoyer un PDF (gère le bruit PPT→PDF)"""
    with pdfplumber.open(filepath) as pdf:
        pages_text = []
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages_text.append(t)
    
    full_text = "\n\n".join(pages_text)
    lines = full_text.split("\n")
    cleaned = []
    
    for line in lines:
        stripped = line.strip()
        # Skip bruit PowerPoint
        if re.match(r'^[a-zéèêëàâäùûüîïôöç]\s*$', stripped):
            continue
        if '© Tous droits réservés' in stripped or 'All rights reserved' in stripped:
            continue
        if 'EM Normandie Business School' in stripped:
            continue
        if re.match(r'^(ip|rto|ze|id|dd|lsn|lsna)$', stripped):
            continue
        if len(stripped) <= 2 and not stripped.isdigit():
            continue
        if 'représentation ou reproduction' in stripped:
            continue
        if stripped in ('PGE', 'Master 1'):
            continue
        if re.match(r'^\d{1,3}$', stripped):
            continue
        if stripped:
            cleaned.append(stripped)
    
    return "\n".join(cleaned)

def extract_and_clean_docx(filepath):
    """Extraire et nettoyer un DOCX"""
    result = subprocess.run(
        ["pandoc", filepath, "-t", "markdown", "--wrap=none"],
        capture_output=True, text=True, timeout=60
    )
    return result.stdout

def extract_and_clean_pptx(filepath):
    """Extraire et nettoyer un PPTX/PPT"""
    ext = Path(filepath).suffix.lower()

    # Pour les .ppt (ancien format), essayer markitdown
    if ext == '.ppt':
        try:
            result = subprocess.run(
                ["python3", "-m", "markitdown", filepath],
                capture_output=True, text=True, timeout=60
            )
            if result.stdout.strip():
                return result.stdout
        except:
            pass
        return f"[Fichier .ppt ancien format - contenu non extractible: {os.path.basename(filepath)}]"

    # Pour les .pptx, utiliser python-pptx directement
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_content.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_content.append(row_text)
            if slide_content:
                slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_content))

        text = "\n\n".join(slides_text)
        # Nettoyage
        lines = text.split("\n")
        cleaned = []
        for line in lines:
            stripped = line.strip()
            if '© Tous droits' in stripped or 'All rights reserved' in stripped:
                continue
            if 'EM Normandie' in stripped:
                continue
            if stripped:
                cleaned.append(line)
        return "\n".join(cleaned)
    except Exception as e:
        return f"[Erreur extraction PPTX: {e}]"

def extract_xlsx(filepath):
    """Extraire les données d'un Excel"""
    wb = openpyxl.load_workbook(filepath, data_only=True)
    text = f"Fichier Excel: {os.path.basename(filepath)}\nFeuilles: {wb.sheetnames}\n\n"
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text += f"=== Feuille: {sheet_name} ===\n"
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            if any(v for v in vals):
                text += " | ".join(vals) + "\n"
        text += "\n"
    return text

def extract_file(filepath):
    """Router l'extraction selon le type"""
    ext = Path(filepath).suffix.lower()
    extractors = {
        '.pdf': extract_and_clean_pdf,
        '.docx': extract_and_clean_docx,
        '.doc': extract_and_clean_docx,
        '.pptx': extract_and_clean_pptx,
        '.ppt': extract_and_clean_pptx,
        '.xlsx': extract_xlsx,
    }
    extractor = extractors.get(ext)
    if extractor:
        return extractor(filepath)
    return open(filepath, 'r', encoding='utf-8', errors='replace').read()

def extract_matiere(matiere_dir, output_dir):
    """Extraire tous les fichiers d'une matière"""
    os.makedirs(output_dir, exist_ok=True)
    results = []
    
    for f in sorted(os.listdir(matiere_dir)):
        filepath = os.path.join(matiere_dir, f)
        if not os.path.isfile(filepath):
            continue
        try:
            content = extract_file(filepath)
            safe_name = re.sub(r'[^\w\-]', '_', os.path.splitext(f)[0])
            out_path = os.path.join(output_dir, f"{safe_name}.txt")
            with open(out_path, 'w', encoding='utf-8') as fout:
                fout.write(content)
            results.append({"filename": f, "extracted": out_path, "chars": len(content)})
            print(f"  ✅ {f} → {len(content):,} chars")
        except Exception as e:
            print(f"  ❌ {f} → {e}")
            results.append({"filename": f, "error": str(e)})
    
    return results

if __name__ == "__main__":
    import sys
    if len(sys.argv) >= 3:
        matiere_dir = sys.argv[1]
        output_dir = sys.argv[2]
        results = extract_matiere(matiere_dir, output_dir)
        print(f"\n✅ {len([r for r in results if 'error' not in r])}/{len(results)} fichiers extraits")
